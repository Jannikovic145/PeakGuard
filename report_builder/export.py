# report_builder/export.py
"""
Export-Funktionen für PeakGuard.
Unterstützt Excel (xlsx) und PowerPoint (pptx) Export.
"""
from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

import pandas as pd

logger = logging.getLogger("peakguard.export")


# ============================================================================
# DATENMODELLE
# ============================================================================
@dataclass
class ExportConfig:
    """Konfiguration für Export"""
    include_raw_data: bool = True
    include_analysis: bool = True
    include_charts: bool = True
    include_recommendations: bool = True
    company_name: str = ""
    site_name: str = ""
    report_date: str = ""


# ============================================================================
# EXCEL EXPORT
# ============================================================================
def export_to_excel(
    df: pd.DataFrame,
    analysis_results: Dict[str, Any],
    output_path: Path,
    config: Optional[ExportConfig] = None,
    progress_callback: Optional[Callable[[float, str], None]] = None
) -> Path:
    """
    Exportiert Analyse-Ergebnisse nach Excel.

    Args:
        df: Rohdaten DataFrame
        analysis_results: Dict mit Analyse-Ergebnissen
        output_path: Ausgabepfad für Excel-Datei
        config: Export-Konfiguration
        progress_callback: Callback für Fortschrittsanzeige

    Returns:
        Pfad zur erstellten Excel-Datei
    """
    try:
        import openpyxl
        from openpyxl.chart import LineChart, Reference
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
        from openpyxl.utils.dataframe import dataframe_to_rows
    except ImportError:
        logger.error("openpyxl nicht installiert. Installieren mit: pip install openpyxl")
        raise ImportError("openpyxl erforderlich für Excel-Export")

    config = config or ExportConfig()

    def report_progress(pct: float, msg: str):
        if progress_callback:
            progress_callback(pct, msg)

    report_progress(0.0, "Excel-Export gestartet...")

    # Workbook erstellen
    wb = openpyxl.Workbook()

    # Styles definieren
    header_fill = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)
    border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )

    # ==================== Sheet 1: Übersicht ====================
    ws_overview = wb.active
    ws_overview.title = "Übersicht"

    overview_data = [
        ["PeakGuard Analyse-Report", ""],
        ["", ""],
        ["Erstellt am:", config.report_date or pd.Timestamp.now().strftime("%d.%m.%Y")],
        ["Unternehmen:", config.company_name or "—"],
        ["Standort:", config.site_name or "—"],
        ["", ""],
        ["Datenzeitraum:", f"{df.index.min()} bis {df.index.max()}" if not df.empty else "—"],
        ["Datenpunkte:", f"{len(df):,}"],
        ["", ""],
    ]

    # KPIs hinzufügen
    if "kpis" in analysis_results:
        overview_data.append(["=== Kennzahlen ===", ""])
        for key, value in analysis_results["kpis"].items():
            overview_data.append([key, str(value)])

    for row_idx, row_data in enumerate(overview_data, 1):
        for col_idx, value in enumerate(row_data, 1):
            cell = ws_overview.cell(row=row_idx, column=col_idx, value=value)
            if row_idx == 1:
                cell.font = Font(bold=True, size=16)
            elif "===" in str(value):
                cell.font = Font(bold=True)

    ws_overview.column_dimensions['A'].width = 25
    ws_overview.column_dimensions['B'].width = 40

    report_progress(0.2, "Übersicht erstellt...")

    # ==================== Sheet 2: Rohdaten ====================
    if config.include_raw_data and not df.empty:
        ws_data = wb.create_sheet("Rohdaten")

        # Header
        df_export = df.reset_index()
        for col_idx, col_name in enumerate(df_export.columns, 1):
            cell = ws_data.cell(row=1, column=col_idx, value=str(col_name))
            cell.fill = header_fill
            cell.font = header_font
            cell.border = border
            cell.alignment = Alignment(horizontal='center')

        # Daten (max 10000 Zeilen für Performance)
        max_rows = min(len(df_export), 10000)
        for row_idx, row in enumerate(df_export.head(max_rows).values, 2):
            for col_idx, value in enumerate(row, 1):
                cell = ws_data.cell(row=row_idx, column=col_idx)
                try:
                    if pd.isna(value):
                        cell.value = ""
                    elif isinstance(value, (pd.Timestamp, datetime)):
                        cell.value = str(value)
                    elif hasattr(value, 'isoformat'):  # datetime-like
                        cell.value = str(value)
                    elif isinstance(value, (int, float)):
                        cell.value = value
                    else:
                        cell.value = str(value)
                except (TypeError, ValueError):
                    cell.value = str(value) if value is not None else ""
                cell.border = border

        # Spaltenbreiten anpassen
        for col in ws_data.columns:
            max_length = 0
            column = col[0].column_letter
            for cell in col[:100]:  # Nur erste 100 Zeilen prüfen
                try:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                except:
                    pass
            ws_data.column_dimensions[column].width = min(max_length + 2, 40)

        if len(df_export) > max_rows:
            ws_data.cell(row=max_rows + 2, column=1,
                        value=f"... {len(df_export) - max_rows} weitere Zeilen nicht exportiert")

    report_progress(0.4, "Rohdaten exportiert...")

    # ==================== Sheet 3: Peak-Analyse ====================
    if config.include_analysis and "peaks" in analysis_results:
        ws_peaks = wb.create_sheet("Peak-Analyse")

        peaks = analysis_results["peaks"]
        peak_headers = ["Rang", "Zeitpunkt", "Leistung (kW)", "Dauer", "Typ"]

        for col_idx, header in enumerate(peak_headers, 1):
            cell = ws_peaks.cell(row=1, column=col_idx, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.border = border

        for row_idx, peak in enumerate(peaks[:20], 2):  # Top 20
            ws_peaks.cell(row=row_idx, column=1, value=row_idx - 1).border = border
            ws_peaks.cell(row=row_idx, column=2, value=str(peak.get("timestamp", ""))).border = border
            ws_peaks.cell(row=row_idx, column=3, value=peak.get("power_kw", 0)).border = border
            ws_peaks.cell(row=row_idx, column=4, value=peak.get("duration", "")).border = border
            ws_peaks.cell(row=row_idx, column=5, value=peak.get("type", "")).border = border

    report_progress(0.6, "Peak-Analyse exportiert...")

    # ==================== Sheet 4: Szenarien ====================
    if config.include_analysis and "scenarios" in analysis_results:
        ws_scenarios = wb.create_sheet("Szenarien")

        scenario_headers = ["Szenario", "Ziel-Cap (kW)", "Einsparung (€/Jahr)", "Reduktion (%)"]

        for col_idx, header in enumerate(scenario_headers, 1):
            cell = ws_scenarios.cell(row=1, column=col_idx, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.border = border

        for row_idx, scenario in enumerate(analysis_results["scenarios"], 2):
            ws_scenarios.cell(row=row_idx, column=1, value=scenario.get("name", "")).border = border
            ws_scenarios.cell(row=row_idx, column=2, value=scenario.get("target_cap", 0)).border = border
            ws_scenarios.cell(row=row_idx, column=3, value=scenario.get("savings", 0)).border = border
            ws_scenarios.cell(row=row_idx, column=4, value=scenario.get("reduction_pct", 0)).border = border

    report_progress(0.8, "Szenarien exportiert...")

    # ==================== Sheet 5: Empfehlungen ====================
    if config.include_recommendations and "recommendations" in analysis_results:
        ws_rec = wb.create_sheet("Empfehlungen")

        rec_headers = ["Priorität", "Kategorie", "Empfehlung", "Einsparpotenzial"]

        for col_idx, header in enumerate(rec_headers, 1):
            cell = ws_rec.cell(row=1, column=col_idx, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.border = border

        for row_idx, rec in enumerate(analysis_results["recommendations"], 2):
            ws_rec.cell(row=row_idx, column=1, value=rec.get("priority", "")).border = border
            ws_rec.cell(row=row_idx, column=2, value=rec.get("category", "")).border = border
            ws_rec.cell(row=row_idx, column=3, value=rec.get("text", "")).border = border
            ws_rec.cell(row=row_idx, column=4, value=rec.get("savings", "")).border = border

        ws_rec.column_dimensions['C'].width = 60

    report_progress(0.9, "Empfehlungen exportiert...")

    # Speichern
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(output_path)

    report_progress(1.0, "Excel-Export abgeschlossen")
    logger.info(f"Excel-Export erfolgreich: {output_path}")

    return output_path


def export_to_excel_bytes(
    df: pd.DataFrame,
    analysis_results: Dict[str, Any],
    config: Optional[ExportConfig] = None
) -> bytes:
    """
    Exportiert nach Excel und gibt Bytes zurück (für Streamlit Download).

    Returns:
        Excel-Datei als Bytes
    """
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
        export_to_excel(df, analysis_results, Path(tmp.name), config)
        tmp.seek(0)
        return Path(tmp.name).read_bytes()


# ============================================================================
# POWERPOINT EXPORT
# ============================================================================
def export_to_powerpoint(
    analysis_results: Dict[str, Any],
    chart_images: Dict[str, bytes],
    output_path: Path,
    config: Optional[ExportConfig] = None,
    progress_callback: Optional[Callable[[float, str], None]] = None
) -> Path:
    """
    Exportiert Analyse-Ergebnisse nach PowerPoint.

    Args:
        analysis_results: Dict mit Analyse-Ergebnissen
        chart_images: Dict mit Chart-Namen und PNG-Bytes
        output_path: Ausgabepfad für PPTX-Datei
        config: Export-Konfiguration
        progress_callback: Callback für Fortschrittsanzeige

    Returns:
        Pfad zur erstellten PowerPoint-Datei
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError:
        logger.error("python-pptx nicht installiert. Installieren mit: pip install python-pptx")
        raise ImportError("python-pptx erforderlich für PowerPoint-Export")

    config = config or ExportConfig()

    def report_progress(pct: float, msg: str):
        if progress_callback:
            progress_callback(pct, msg)

    report_progress(0.0, "PowerPoint-Export gestartet...")

    # Präsentation erstellen
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    # Farben
    primary_color = RGBColor(0x1E, 0x3A, 0x5F)  # PeakGuard Dunkelblau
    accent_color = RGBColor(0xE6, 0x7E, 0x22)   # Orange

    # ==================== Folie 1: Titelfolie ====================
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "PeakGuard Lastanalyse"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER

    # Untertitel
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = f"{config.site_name or 'Standort'} | {config.report_date or 'Datum'}"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    subtitle_para.alignment = PP_ALIGN.CENTER

    report_progress(0.2, "Titelfolie erstellt...")

    # ==================== Folie 2: Executive Summary ====================
    slide = prs.slides.add_slide(slide_layout)

    # Überschrift
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    header_frame = header_box.text_frame
    header_para = header_frame.paragraphs[0]
    header_para.text = "Executive Summary"
    header_para.font.size = Pt(32)
    header_para.font.bold = True
    header_para.font.color.rgb = primary_color

    # KPIs als Karten
    if "kpis" in analysis_results:
        kpis = analysis_results["kpis"]
        kpi_items = list(kpis.items())[:6]  # Max 6 KPIs

        for idx, (key, value) in enumerate(kpi_items):
            row = idx // 3
            col = idx % 3
            x = Inches(0.5 + col * 4.2)
            y = Inches(1.5 + row * 2.5)

            # KPI Box
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                x, y, Inches(4), Inches(2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

            # KPI Wert
            value_box = slide.shapes.add_textbox(x, y + Inches(0.3), Inches(4), Inches(1))
            value_frame = value_box.text_frame
            value_para = value_frame.paragraphs[0]
            value_para.text = str(value)
            value_para.font.size = Pt(28)
            value_para.font.bold = True
            value_para.font.color.rgb = primary_color
            value_para.alignment = PP_ALIGN.CENTER

            # KPI Label
            label_box = slide.shapes.add_textbox(x, y + Inches(1.2), Inches(4), Inches(0.5))
            label_frame = label_box.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = key
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            label_para.alignment = PP_ALIGN.CENTER

    report_progress(0.4, "Summary erstellt...")

    # ==================== Folien 3+: Charts ====================
    chart_titles = {
        "timeseries": "Lastgang-Übersicht",
        "duration_curve": "Jahresdauerlinie",
        "heatmap": "Lastheatmap",
        "monthly_peaks": "Monatliche Spitzenlasten",
        "peak_events": "Peak-Events-Analyse"
    }

    for idx, (chart_name, chart_bytes) in enumerate(chart_images.items()):
        slide = prs.slides.add_slide(slide_layout)

        # Überschrift
        title = chart_titles.get(chart_name, chart_name)
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        header_frame = header_box.text_frame
        header_para = header_frame.paragraphs[0]
        header_para.text = title
        header_para.font.size = Pt(28)
        header_para.font.bold = True
        header_para.font.color.rgb = primary_color

        # Chart-Bild einfügen
        if chart_bytes:
            image_stream = io.BytesIO(chart_bytes)
            slide.shapes.add_picture(
                image_stream,
                Inches(0.5), Inches(1.2),
                width=Inches(12)
            )

        report_progress(0.4 + (idx + 1) * 0.1, f"Chart {chart_name} eingefügt...")

    # ==================== Letzte Folie: Empfehlungen ====================
    if "recommendations" in analysis_results:
        slide = prs.slides.add_slide(slide_layout)

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        header_frame = header_box.text_frame
        header_para = header_frame.paragraphs[0]
        header_para.text = "Handlungsempfehlungen"
        header_para.font.size = Pt(28)
        header_para.font.bold = True
        header_para.font.color.rgb = primary_color

        rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(6))
        rec_frame = rec_box.text_frame
        rec_frame.word_wrap = True

        for rec in analysis_results["recommendations"][:5]:  # Top 5
            para = rec_frame.add_paragraph()
            para.text = f"• {rec.get('text', '')}"
            para.font.size = Pt(16)
            para.space_after = Pt(12)

            if rec.get("savings"):
                para2 = rec_frame.add_paragraph()
                para2.text = f"   Potenzial: {rec.get('savings')}"
                para2.font.size = Pt(14)
                para2.font.color.rgb = accent_color
                para2.space_after = Pt(18)

    report_progress(0.9, "Empfehlungen eingefügt...")

    # Speichern
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

    report_progress(1.0, "PowerPoint-Export abgeschlossen")
    logger.info(f"PowerPoint-Export erfolgreich: {output_path}")

    return output_path


def export_to_powerpoint_bytes(
    analysis_results: Dict[str, Any],
    chart_images: Dict[str, bytes],
    config: Optional[ExportConfig] = None
) -> bytes:
    """
    Exportiert nach PowerPoint und gibt Bytes zurück (für Streamlit Download).

    Returns:
        PPTX-Datei als Bytes
    """
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        export_to_powerpoint(analysis_results, chart_images, Path(tmp.name), config)
        return Path(tmp.name).read_bytes()


# ============================================================================
# CSV EXPORT (VEREINFACHT)
# ============================================================================
def export_analysis_to_csv(
    df: pd.DataFrame,
    output_path: Path,
    include_index: bool = True
) -> Path:
    """
    Exportiert DataFrame nach CSV mit deutscher Formatierung.

    Args:
        df: DataFrame
        output_path: Ausgabepfad
        include_index: Index mit exportieren

    Returns:
        Pfad zur CSV-Datei
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    df.to_csv(
        output_path,
        sep=";",
        decimal=",",
        index=include_index,
        encoding="utf-8-sig"
    )

    logger.info(f"CSV-Export erfolgreich: {output_path}")
    return output_path
